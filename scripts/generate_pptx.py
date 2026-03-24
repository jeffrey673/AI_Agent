#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Generate SKIN1004 AI Agent presentation as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
ORANGE = RGBColor(0xE8, 0x92, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
DARKER = RGBColor(0x0F, 0x0F, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARKER):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()


def add_slide_num(slide, num, total):
    txBox = slide.shapes.add_textbox(W - Inches(1.5), H - Inches(0.5), Inches(1.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, text, y=Inches(0.8), size=Pt(36)):
    txBox = slide.shapes.add_textbox(Inches(1), y, Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = WHITE
    return txBox


def add_accent_title(slide, normal, accent, y=Inches(0.8)):
    txBox = slide.shapes.add_textbox(Inches(1), y, Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = normal
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.color.rgb = WHITE
    run2 = p.add_run()
    run2.text = accent
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = ORANGE


def add_body(slide, text, x=Inches(1), y=Inches(1.8), w=Inches(11), size=Pt(16), color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(x, y, w, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def add_table(slide, headers, rows, x=Inches(1), y=Inches(2), w=Inches(11.3)):
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(len(rows) + 1, cols, x, y, w, Inches(0.4 * (len(rows) + 1)))
    tbl = tbl_shape.table

    # Header
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x3D, 0x2E, 0x00)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ORANGE

    # Rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARKER
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = LIGHT_GRAY

    return tbl_shape


def add_card(slide, emoji, label, desc, x, y, w=Inches(3.5), h=Inches(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x38)
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x55)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(12)

    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(28)

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE

    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(13)
    p3.font.color.rgb = LIGHT_GRAY


def add_stat(slide, num, desc, x, y, w=Inches(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x20, 0x00)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(16)

    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER


TOTAL = 15

# =============================================
# Slide 1: Title
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)
add_top_bar(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "SKIN1004"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Enterprise AI Agent"
p2.font.size = Pt(44)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.5))
p = sub.text_frame.paragraphs[0]
p.text = '"What Do You Crave? — 데이터를 묻고, 답을 얻다"'
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

meta = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(1))
tf = meta.text_frame
for line in ["Craver Corporation", "발표: 임재필 (Jeffrey Im) | 데이터분석파트", "2026년 3월"]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
add_slide_num(slide, 1, TOTAL)

# =============================================
# Slide 2: Why
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "왜 ", "만들었는가?")

add_body(slide, "❌ 기존 업무의 문제점", y=Inches(1.7), size=Pt(18), color=ORANGE)
add_body(slide,
    "• 매출 확인: BigQuery 콘솔 → SQL 작성 → 해석 (20분+)\n"
    "• 문서 검색: Notion 페이지 뒤지기\n"
    "• CS 응대: 스프레드시트 13탭 수동 검색\n"
    "• 메일/일정: 앱 전환 필요",
    y=Inches(2.2), size=Pt(15))

add_body(slide, "✅ AI Agent 도입 후", x=Inches(7), y=Inches(1.7), size=Pt(18), color=ORANGE, w=Inches(5))
add_body(slide,
    '"이번 달 인도네시아 매출 알려줘"\n→ 10초 내 답변 + 차트 자동 생성\n\n'
    "• 자연어로 모든 업무 데이터에 즉시 접근\n"
    "• SQL 몰라도 데이터 분석 가능\n"
    "• 전직원 사용 가능 (AD 연동)",
    x=Inches(7), y=Inches(2.2), w=Inches(5), size=Pt(15))
add_slide_num(slide, 2, TOTAL)

# =============================================
# Slide 3: Architecture
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "시스템 ", "아키텍처")

# User
add_card(slide, "👤", "사용자", "웹 브라우저", Inches(5), Inches(1.5), Inches(3.3), Inches(1))
# Server
add_card(slide, "⚡", "FastAPI + LangGraph", "질문 의도 자동 분류 (6개 경로)", Inches(3.5), Inches(3), Inches(6.3), Inches(1.2))
# 6 routes
labels = [("📊","BQ\n매출"), ("📋","Notion\n문서"), ("🧴","CS\nQ&A"), ("📧","GWS\n메일"), ("📈","Multi\n종합"), ("💬","Direct\n대화")]
for i, (emoji, lbl) in enumerate(labels):
    add_card(slide, emoji, lbl, "", Inches(0.8 + i * 2.05), Inches(4.8), Inches(1.85), Inches(1.5))

add_slide_num(slide, 3, TOTAL)

# =============================================
# Slide 4: 6 Routes
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "6가지 ", "자동 라우팅")
add_table(slide,
    ["라우트", "질문 예시", "처리 방식"],
    [
        ["📊 BigQuery", '"이번 달 매출", "국가별 TOP 5"', "자연어→SQL→실행→차트"],
        ["📋 Notion", '"틱톡샵 접속 방법", "출장 가이드"', "사내 문서 벡터 검색"],
        ["🧴 CS Q&A", '"센텔라 앰플 성분", "사용법"', "739개 CS 키워드 검색"],
        ["📧 GWS", '"오늘 일정", "메일 확인"', "Gmail/Calendar/Drive API"],
        ["📈 Multi", '"매출 하락 원인 분석"', "BQ + 웹검색 종합"],
        ["💬 Direct", '"안녕", "날씨", "회사 소개"', "Claude Sonnet 실시간"],
    ])
add_slide_num(slide, 4, TOTAL)

# =============================================
# Slide 5: Tech Stack
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "기술 ", "스택")

cards = [
    ("🤖", "LLM (대화)", "Claude Sonnet 4\nTTFB 2초"),
    ("⚡", "LLM (SQL)", "Gemini 2.5 Flash\nSQL + 차트"),
    ("🔄", "오케스트레이션", "LangGraph\nAI 워크플로우"),
    ("📊", "데이터베이스", "BigQuery\n13개 테이블"),
    ("🔐", "인증", "AD + JWT\n그룹 권한"),
    ("🌐", "프론트엔드", "자체 SPA\nSSE 스트리밍"),
]
for i, (emoji, label, desc) in enumerate(cards):
    col = i % 3
    row = i // 3
    add_card(slide, emoji, label, desc, Inches(0.8 + col * 4), Inches(1.8 + row * 2.5), Inches(3.5), Inches(2))
add_slide_num(slide, 5, TOTAL)

# =============================================
# Slide 6: 13 Data Sources
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "13개 ", "데이터 소스")

sources = [
    ("📊","통합 매출"), ("📦","제품"), ("💰","광고비"), ("💵","마케팅비용"),
    ("👤","인플루언서"), ("🔍","아마존검색"), ("📱","메타광고"), ("🛒","플랫폼순위"),
    ("🛍️","Shopify"), ("⭐","리뷰 4종"), ("📋","Notion"), ("🧴","CS Q&A"),
]
for i, (emoji, label) in enumerate(sources):
    col = i % 4
    row = i // 4
    add_card(slide, emoji, label, "", Inches(0.6 + col * 3.1), Inches(1.7 + row * 1.8), Inches(2.8), Inches(1.4))
add_slide_num(slide, 6, TOTAL)

# =============================================
# Slide 7: Streaming
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "실시간 ", "스트리밍")

add_stat(slide, "9.5초", "이전 — 전체 완성 후 표시", Inches(1.5), Inches(2), Inches(4))
add_stat(slide, "2.0초", "현재 — 첫 텍스트 즉시 타이핑", Inches(7), Inches(2), Inches(4.5))

add_body(slide,
    "→  ChatGPT급 체감 속도\n"
    "→  블링킹 커서 + 50ms 렌더링\n"
    "→  Claude Sonnet 스트리밍 → Thread + Queue → SSE",
    y=Inches(4.2), size=Pt(16))
add_slide_num(slide, 7, TOTAL)

# =============================================
# Slide 8: Speed
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "속도 ", "최적화")
add_table(slide,
    ["항목", "이전 (v8.2)", "이후 (v9.0)", "개선"],
    [
        ["일반 대화 (TTFB)", "9.5초", "2.0초", "-79%"],
        ["BQ 간단 쿼리", "18.1초", "9.5초", "-47%"],
        ["BQ 캐시 히트", "18.1초", "7.5초", "-59%"],
        ["웹검색 질문", "17.7초", "7.2초", "-59%"],
    ])
add_stat(slide, "86%", "속도 개선", Inches(1), Inches(5.2), Inches(3.5))
add_stat(slide, "0건", "Timeout", Inches(5), Inches(5.2), Inches(3.5))
add_stat(slide, "$400", "월 운영비", Inches(9), Inches(5.2), Inches(3.5))
add_slide_num(slide, 8, TOTAL)

# =============================================
# Slide 9: Security
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "보안 & ", "권한 관리")

add_body(slide, "🔐 인증 체계", y=Inches(1.7), size=Pt(18), color=ORANGE)
add_body(slide,
    "• AD 연동 — 328명 전직원 자동 등록\n"
    "• JWT httpOnly 쿠키 인증\n"
    "• GWS: 개인별 Google OAuth (타인 접근 불가)",
    y=Inches(2.2), size=Pt(15))

add_body(slide, "🛡️ 데이터 보안", x=Inches(7), y=Inches(1.7), size=Pt(18), color=ORANGE, w=Inches(5))
add_body(slide,
    "• 그룹별 데이터 접근 제한 (brand_filter)\n"
    "  SK그룹: SK/CL/CBT | DD그룹: DD만\n"
    "• SQL 안전장치: SELECT ONLY\n"
    "• 사내 서버 운영 (외부 전송 없음)",
    x=Inches(7), y=Inches(2.2), w=Inches(5), size=Pt(15))
add_slide_num(slide, 9, TOTAL)

# =============================================
# Slide 10: QA
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "QA ", "테스트 결과")

add_stat(slide, "99.6%", "API QA PASS\n(900문항)", Inches(1), Inches(2), Inches(3.5))
add_stat(slide, "93%", "Playwright E2E\n(브라우저 100문항)", Inches(5), Inches(2), Inches(3.5))
add_stat(slide, "77%", "jp2 이슈 해결\n(26건 중 20건)", Inches(9), Inches(2), Inches(3.5))

add_body(slide,
    "9개 테이블 × 100문항 = 900건 자동 테스트\n"
    "Playwright로 실제 브라우저에서 로그인 → 질문 → 답변 검증",
    y=Inches(4.5), size=Pt(16))
add_slide_num(slide, 10, TOTAL)

# =============================================
# Slide 11: vs ChatGPT
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "vs ", "ChatGPT")
add_table(slide,
    ["기능", "ChatGPT", "SKIN1004 AI Agent"],
    [
        ["매출 데이터 조회", "❌ 불가", "✅ BigQuery 실시간"],
        ["사내 문서 검색", "❌ 불가", "✅ Notion 연동"],
        ["CS 제품 Q&A", "❌ 불가", "✅ 739개 DB"],
        ["메일/일정", "❌ 불가", "✅ 개인별 OAuth"],
        ["차트 생성", "⚠️ 제한적", "✅ Chart.js 인터랙티브"],
        ["데이터 보안", "⚠️ 외부 전송", "✅ 사내 서버"],
        ["스트리밍", "✅", "✅ ChatGPT급"],
    ])
add_slide_num(slide, 11, TOTAL)

# =============================================
# Slide 12: Guide
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "사용 ", "가이드")

add_body(slide,
    "1️⃣  로그인: 부서 → 이름 → 비밀번호\n"
    "2️⃣  질문: 자연어로 무엇이든 입력\n"
    "3️⃣  결과: 표, 차트, 인사이트 자동 생성\n"
    "4️⃣  후속: 대화 맥락 유지 심화 분석",
    y=Inches(1.8), size=Pt(17))

add_table(slide,
    ["부서", "질문 예시"],
    [
        ["B2B팀", '"거래처별 매출 순위"'],
        ["GM팀", '"인도네시아 쇼피 월별 추이"'],
        ["마케팅", '"Facebook ROAS 국가별 비교"'],
        ["경영관리", '"전사 1분기 실적 요약"'],
        ["CS", '"센텔라 앰플 성분 알려줘"'],
    ], y=Inches(3.8))
add_slide_num(slide, 12, TOTAL)

# =============================================
# Slide 13: Advanced
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "고급 ", "기능")

features = [
    ("📈", "자동 차트", "바/라인/파이 차트\n자동 생성"),
    ("🔍", "웹 검색", "날씨, 뉴스, 환율\nGoogle 실시간"),
    ("📧", "Google 연동", "메일, 일정, 드라이브\n자동 OAuth"),
    ("🖼️", "이미지 분석", "업로드 → AI 분석"),
    ("🌍", "다국어", "영어/스페인어\n해당 언어 답변"),
    ("⏹️", "생성 중지", "Stop 버튼\n응답 즉시 중단"),
]
for i, (emoji, label, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    add_card(slide, emoji, label, desc, Inches(0.8 + col * 4), Inches(1.8 + row * 2.5), Inches(3.5), Inches(2))
add_slide_num(slide, 13, TOTAL)

# =============================================
# Slide 14: Roadmap
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)
add_accent_title(slide, "향후 ", "로드맵")

add_body(slide,
    "🟠 2026 Q2  자동 일일 리포트 → Slack/Teams 전송\n"
    "🟠 2026 Q2  음성 입력/출력 (STT/TTS)\n"
    "🔵 2026 Q3  모바일 최적화 (PWA)\n"
    "🔵 2026 Q3  예측 분석 — 매출/재고 예측\n"
    "🟢 2026 Q4  다국어 CS 자동 응대\n"
    "🟢 2027     전사 업무 자동화 플랫폼",
    y=Inches(1.8), size=Pt(18))

add_body(slide,
    '비전: "모든 크레이버 직원이 데이터 분석가가 되는 세상"',
    y=Inches(5.5), size=Pt(20), color=ORANGE)
add_slide_num(slide, 14, TOTAL)

# =============================================
# Slide 15: Q&A
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Q&A"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.5))
p = sub.text_frame.paragraphs[0]
p.text = '"무엇이 궁금하신가요?"'
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

meta = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(2))
tf = meta.text_frame
for line, color in [
    ("🎯 실시간 데모 가능", LIGHT_GRAY),
    ("", GRAY),
    ("임재필 (Jeffrey Im)", WHITE),
    ("jeffrey@skin1004korea.com", GRAY),
    ("", GRAY),
    ("What Do You Crave?", ORANGE),
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
add_slide_num(slide, 15, TOTAL)

# Save
output = "docs/SKIN1004_AI_Agent_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
